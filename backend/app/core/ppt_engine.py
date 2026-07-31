import os
import io
import re
import logging
import base64
import shutil
import subprocess
import tempfile
import textwrap
from typing import Dict, Any, List, Optional
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE_TYPE

from reportlab.lib.pagesizes import letter, landscape
from reportlab.lib import colors
from reportlab.pdfgen import canvas
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak, KeepTogether, HRFlowable
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle

logger = logging.getLogger("kairos.ppt_engine")

# Resolve templates directory relative to this file's location.
# Works on local dev, Docker (/home/user/app), and Render (/opt/render/project/src).
_THIS_DIR = os.path.dirname(os.path.abspath(__file__))
TEMPLATES_DIR = os.path.join(os.path.dirname(_THIS_DIR), "static", "ppt_templates")

# Log template status at import time so deployment issues are immediately visible.
if os.path.isdir(TEMPLATES_DIR):
    _found = [f for f in os.listdir(TEMPLATES_DIR) if f.endswith(".pptx")]
    logger.info("PPT templates dir: %s  (%d templates found: %s)", TEMPLATES_DIR, len(_found), _found)
else:
    logger.warning("PPT templates dir NOT FOUND at %s — presentation exports will fail!", TEMPLATES_DIR)

PREDEFINED_TEMPLATES = {
    "template-1": {
        "id": "template-1",
        "name": "Cyber Neon Executive",
        "file": "Template-1.pptx",
        "description": "High-impact dark theme with neon purple accents. Ideal for AI & Tech Hackathons.",
        "slides_count": 10
    },
    "template-2": {
        "id": "template-2",
        "name": "Minimalist Modern Tech",
        "file": "Template-2.pptx",
        "description": "Clean, structured slide deck with high legibility and sleek grid layouts.",
        "slides_count": 11
    },
    "template-3": {
        "id": "template-3",
        "name": "Vibrant Launchpad",
        "file": "Template-3.pptx",
        "description": "Dynamic layout designed for pitch competitions, featuring prominent metric boxes.",
        "slides_count": 11
    },
    "template-4": {
        "id": "template-4",
        "name": "Enterprise Architecture",
        "file": "Template-4.pptx",
        "description": "Comprehensive design focusing on workflow, technical specs, and milestone flowcharts.",
        "slides_count": 13
    },
    "template-5": {
        "id": "template-5",
        "name": "Futuristic AI Studio",
        "file": "Template-5.pptx",
        "description": "Sleek obsidian gradient deck highlighting AI solution capabilities and team power.",
        "slides_count": 10
    }
}

class NumberedCanvas(canvas.Canvas):
    def __init__(self, *args, **kwargs):
        super().__init__(*args, **kwargs)
        self._saved_page_states = []

    def showPage(self):
        self._saved_page_states.append(dict(self.__dict__))
        self._startPage()

    def save(self):
        num_pages = len(self._saved_page_states)
        for state in self._saved_page_states:
            self.__dict__.update(state)
            self.draw_page_decorations(num_pages)
            super().showPage()
        super().save()

    def draw_page_decorations(self, page_count):
        self.saveState()
        
        # Tactical Header Banner on pages after cover
        if self._pageNumber > 1:
            self.setStrokeColor(colors.HexColor("#0f172a"))
            self.setLineWidth(1.2)
            self.line(36, 756, 576, 756)
            self.setFont("Helvetica-Bold", 8)
            self.setFillColor(colors.HexColor("#dc2626"))
            self.drawString(36, 762, "[ RESTRICTED // MIL-SPEC OPERATIONAL BLUEPRINT ]")
            self.setFont("Helvetica-Bold", 8)
            self.setFillColor(colors.HexColor("#475569"))
            self.drawRightString(576, 762, "SECURITY CLEARANCE: LEVEL 5")

        # Military Footer Line & Page Numbers
        self.setStrokeColor(colors.HexColor("#334155"))
        self.setLineWidth(1)
        self.line(36, 40, 576, 40)
        
        self.setFont("Helvetica-Bold", 8)
        self.setFillColor(colors.HexColor("#0f172a"))
        self.drawString(36, 26, "KAIROS TACTICAL ENGINE • OPERATIONAL SPECIFICATIONS DOSSIER")
        
        page_text = f"PAGE {self._pageNumber} OF {page_count}"
        self.setFont("Helvetica-Bold", 8)
        self.setFillColor(colors.HexColor("#6d28d9"))
        self.drawRightString(576, 26, page_text)
        self.restoreState()


class PPTEngine:

    @staticmethod
    def get_template_path(template_id: str) -> str:
        t_info = PREDEFINED_TEMPLATES.get(template_id, PREDEFINED_TEMPLATES["template-1"])
        return os.path.join(TEMPLATES_DIR, t_info["file"])

    @staticmethod
    def analyze_presentation(prs: Presentation) -> Dict[str, Any]:
        slide_details = []
        total_shapes = 0
        total_text_frames = 0

        for idx, slide in enumerate(prs.slides):
            text_boxes = []
            for shape in slide.shapes:
                total_shapes += 1
                if shape.has_text_frame:
                    total_text_frames += 1
                    text_content = shape.text_frame.text.strip()
                    if text_content:
                        text_boxes.append({
                            "shape_name": shape.name,
                            "char_count": len(text_content),
                            "sample_text": text_content[:60]
                        })
            
            slide_details.append({
                "slide_number": idx + 1,
                "shape_count": len(slide.shapes),
                "text_box_count": len(text_boxes),
                "sample_content": text_boxes[:3]
            })

        return {
            "slide_count": len(prs.slides),
            "total_shapes": total_shapes,
            "total_text_frames": total_text_frames,
            "slide_width_inches": round(prs.slide_width.inches, 2),
            "slide_height_inches": round(prs.slide_height.inches, 2),
            "slide_details": slide_details
        }

    @staticmethod
    def fit_text_to_frame(text_frame, text: str, max_font_size: int = 16, min_font_size: int = 9, is_title: bool = False):
        """
        Replace a template's copy without changing its geometry or visual style.

        Template text boxes are deliberately small.  Keeping their original font
        size for arbitrary LLM output is what caused the old exporter to produce
        overflowing and clipped slides, so this method measures an approximate
        line budget and reduces the inherited size only when it is necessary.
        """
        text_frame.word_wrap = True
        cleaned = re.sub(r"[ \t]+", " ", str(text or "")).strip()
        if not cleaned:
            text_frame.clear()
            return

        existing_paras = list(text_frame.paragraphs)
        if len(existing_paras) == 0:
            return

        reference_para = existing_paras[0]
        reference_run = reference_para.runs[0] if reference_para.runs else None
        original_pt = (
            reference_run.font.size.pt
            if reference_run and reference_run.font.size
            else float(max_font_size)
        )
        min_pt = max(float(min_font_size), 7.0)

        # A conservative approximation is preferable to PowerPoint clipping
        # text after export.  PowerPoint's actual font metrics vary by machine,
        # so leave a little breathing room in both dimensions.
        width_pt = 360.0
        shape = getattr(text_frame, "_parent", None)
        if shape is not None:
            margin_x = (float(text_frame.margin_left) + float(text_frame.margin_right)) / 12700
            margin_y = (float(text_frame.margin_top) + float(text_frame.margin_bottom)) / 12700
            width_pt = max(24.0, shape.width / 12700 - margin_x)
            height_pt = max(12.0, shape.height / 12700 - margin_y)
        else:
            width_pt, height_pt = 360.0, 120.0

        raw_lines = [line.strip() for line in cleaned.splitlines() if line.strip()] or [cleaned]
        bullet_lines = []
        for line in raw_lines:
            is_bullet = line.startswith(("- ", "* ", "• "))
            value = line[2:].strip() if is_bullet else line
            bullet_lines.append(("• " if is_bullet else "") + value)

        def wrap_for_size(font_pt: float):
            chars_per_line = max(8, int((width_pt / max(font_pt, 1)) * 1.95))
            wrapped = []
            for line in bullet_lines:
                prefix = "• " if line.startswith("• ") else ""
                value = line[len(prefix):]
                pieces = textwrap.wrap(
                    value,
                    width=max(8, chars_per_line - len(prefix)),
                    break_long_words=True,
                    break_on_hyphens=False,
                ) or [""]
                wrapped.extend([prefix + pieces[0]] + pieces[1:])
            max_lines = max(1, int(height_pt / max(font_pt * 1.2, 1)))
            return wrapped, max_lines

        chosen_pt = original_pt
        wrapped, max_lines = wrap_for_size(chosen_pt)
        while len(wrapped) > max_lines and chosen_pt > min_pt:
            chosen_pt = max(min_pt, chosen_pt - 1)
            wrapped, max_lines = wrap_for_size(chosen_pt)

        if len(wrapped) > max_lines:
            wrapped = wrapped[:max_lines]
            last = wrapped[-1].rstrip()
            wrapped[-1] = (last[: max(1, len(last) - 3)].rstrip() + "...")

        style = {
            "size": Pt(chosen_pt),
            "name": reference_run.font.name if reference_run else None,
            "bold": reference_run.font.bold if reference_run else None,
            "italic": reference_run.font.italic if reference_run else None,
            "color": None,
        }
        if reference_run:
            try:
                if reference_run.font.color.type == 1:
                    style["color"] = reference_run.font.color.rgb
            except Exception:
                pass

        alignment = reference_para.alignment
        space_before = reference_para.space_before
        space_after = reference_para.space_after
        line_spacing = reference_para.line_spacing
        level = reference_para.level

        text_frame.clear()
        for index, line in enumerate(wrapped):
            paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
            paragraph.alignment = alignment
            paragraph.level = level
            paragraph.space_before = space_before
            paragraph.space_after = space_after
            paragraph.line_spacing = line_spacing
            run = paragraph.add_run()
            run.text = line
            run.font.size = style["size"]
            if style["name"]:
                run.font.name = style["name"]
            if style["bold"] is not None:
                run.font.bold = style["bold"]
            if style["italic"] is not None:
                run.font.italic = style["italic"]
            if style["color"] is not None:
                run.font.color.rgb = style["color"]

        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.NONE

    @staticmethod
    def _rgb_from_color(color, default=(255, 255, 255)):
        try:
            if color and color.type == 1 and color.rgb:
                value = str(color.rgb)
                return tuple(int(value[i:i + 2], 16) for i in (0, 2, 4))
        except Exception:
            pass
        return default

    @staticmethod
    def _shape_fill_color(shape):
        try:
            if shape.fill.type is not None:
                return PPTEngine._rgb_from_color(shape.fill.fore_color, None)
        except Exception:
            pass
        return None

    @staticmethod
    def _shape_line_color(shape):
        try:
            if shape.line.fill.type is not None:
                return PPTEngine._rgb_from_color(shape.line.color, None)
        except Exception:
            pass
        return None

    @staticmethod
    def _load_preview_font(ImageFont, bold=False, size=16):
        candidates = []
        if bold:
            candidates.extend([
                "/usr/share/fonts/truetype/liberation2/LiberationSans-Bold.ttf",
                "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
            ])
        else:
            candidates.extend([
                "/usr/share/fonts/truetype/liberation2/LiberationSans-Regular.ttf",
                "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
            ])
        candidates.append("/System/Library/Fonts/Helvetica.ttc")
        for path in candidates:
            try:
                return ImageFont.truetype(path, max(1, size))
            except Exception:
                continue
        return ImageFont.load_default()

    @staticmethod
    def _render_with_libreoffice(pptx_bytes: bytes, scale: float):
        """Render using PowerPoint-compatible LibreOffice when installed."""
        if not PPTEngine.native_renderer_available():
            return None
        libreoffice = shutil.which("libreoffice") or shutil.which("soffice")
        pdftoppm = shutil.which("pdftoppm")

        with tempfile.TemporaryDirectory(prefix="kairos-ppt-") as tmp:
            source = os.path.join(tmp, "presentation.pptx")
            with open(source, "wb") as handle:
                handle.write(pptx_bytes)
            converted = subprocess.run(
                [libreoffice, "--headless", "--convert-to", "pdf", "--outdir", tmp, source],
                capture_output=True,
                timeout=60,
                check=False,
            )
            pdf_path = os.path.join(tmp, "presentation.pdf")
            if converted.returncode != 0 or not os.path.exists(pdf_path):
                logger.warning("LibreOffice could not render PPTX: %s", converted.stderr[-500:].decode(errors="ignore"))
                return None

            prefix = os.path.join(tmp, "slide")
            rasterized = subprocess.run(
                [pdftoppm, "-png", "-r", str(max(72, int(96 * scale))), pdf_path, prefix],
                capture_output=True,
                timeout=60,
                check=False,
            )
            if rasterized.returncode != 0:
                return None
            paths = sorted(
                (os.path.join(tmp, name) for name in os.listdir(tmp) if name.startswith("slide-") and name.endswith(".png")),
                key=lambda path: int(re.search(r"-(\d+)\.png$", path).group(1)),
            )
            if not paths:
                return None
            encoded = []
            for path in paths:
                with open(path, "rb") as handle:
                    encoded.append(base64.b64encode(handle.read()).decode("ascii"))
            return encoded

    @staticmethod
    def native_renderer_available() -> bool:
        return bool((shutil.which("libreoffice") or shutil.which("soffice")) and shutil.which("pdftoppm"))

    @staticmethod
    def render_slides_as_images(pptx_bytes: bytes, scale: float = 1.5) -> list:
        """
        Render the generated deck for the live preview.

        Native LibreOffice rendering is used in production containers so master
        graphics, gradients, cropping, and fonts match the downloaded deck.  A
        Pillow renderer remains as a deterministic local fallback for machines
        without LibreOffice.
        """
        native = PPTEngine._render_with_libreoffice(pptx_bytes, scale)
        if native:
            return native

        from PIL import Image, ImageDraw, ImageFont, ImageOps

        prs = Presentation(io.BytesIO(pptx_bytes))
        slide_w_px = int(prs.slide_width.inches * 96 * scale)
        slide_h_px = int(prs.slide_height.inches * 96 * scale)
        emu_to_px = lambda emu: int(emu / 914400 * 96 * scale)

        slide_images = []

        for slide in prs.slides:
            img = Image.new("RGB", (slide_w_px, slide_h_px), (255, 255, 255))
            draw = ImageDraw.Draw(img)

            # Check slide background fill color
            bg = slide.background
            if bg and bg.fill and bg.fill.type is not None:
                try:
                    fc = bg.fill.fore_color
                    if fc and fc.type == 1:
                        draw.rectangle([(0, 0), (slide_w_px, slide_h_px)], fill=PPTEngine._rgb_from_color(fc, (255, 255, 255)))
                except Exception:
                    pass

            # Draw in z-order.  The old preview rendered only text and pictures,
            # which made every template's cards, bands, and panels disappear.
            def iter_shapes(container, offset_x=0, offset_y=0):
                for child in container.shapes:
                    yield child, offset_x, offset_y
                    if child.shape_type == MSO_SHAPE_TYPE.GROUP:
                        yield from iter_shapes(
                            child,
                            offset_x + emu_to_px(child.left),
                            offset_y + emu_to_px(child.top),
                        )

            for shape, offset_x, offset_y in iter_shapes(slide):
                x = offset_x + emu_to_px(shape.left)
                y = offset_y + emu_to_px(shape.top)
                w = max(1, emu_to_px(shape.width))
                h = max(1, emu_to_px(shape.height))

                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    continue

                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        pic_img = Image.open(io.BytesIO(shape.image.blob)).convert("RGB")
                        pic_img = ImageOps.fit(pic_img, (w, h), method=Image.Resampling.LANCZOS, centering=(0.5, 0.5))
                        img.paste(pic_img, (x, y))
                    except Exception as e:
                        logger.warning(f"Failed to render slide picture shape in preview: {e}")
                    continue

                fill = PPTEngine._shape_fill_color(shape)
                if fill:
                    try:
                        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and getattr(shape.auto_shape_type, "name", "") in {"OVAL", "ELLIPSE"}:
                            draw.ellipse((x, y, x + w, y + h), fill=fill)
                        else:
                            draw.rectangle((x, y, x + w, y + h), fill=fill)
                    except Exception:
                        draw.rectangle((x, y, x + w, y + h), fill=fill)
                line_color = PPTEngine._shape_line_color(shape)
                if line_color:
                    draw.rectangle((x, y, x + w, y + h), outline=line_color, width=max(1, int(scale)))

                if not shape.has_text_frame:
                    continue
                tf = shape.text_frame
                txt = tf.text.strip()
                if not txt:
                    continue

                # Gather all lines with their font properties
                text_margin_x = max(2, emu_to_px(tf.margin_left))
                text_margin_y = max(2, emu_to_px(tf.margin_top))
                text_y = y + text_margin_y
                rendered_lines = []
                for para in tf.paragraphs:
                    line_text = ""
                    font_size_px = 16
                    font_color = (50, 50, 50)
                    is_bold = False

                    for run in para.runs:
                        if not run.text:
                            continue
                        line_text += run.text

                        # Get font size
                        if run.font.size:
                            raw_pt = run.font.size / 12700
                            capped_pt = min(raw_pt, 72)
                            font_size_px = int(capped_pt * scale)

                        # Get font color
                        try:
                            if run.font.color and run.font.color.type == 1:
                                c = run.font.color.rgb
                                cs = str(c)
                                font_color = (int(cs[0:2], 16), int(cs[2:4], 16), int(cs[4:6], 16))
                        except Exception:
                            pass

                        if run.font.bold:
                            is_bold = True

                    if not line_text.strip():
                        y += font_size_px + 4
                        continue

                    pil_font = PPTEngine._load_preview_font(ImageFont, is_bold, font_size_px)

                    # Word-wrap text to fit within shape width
                    wrapped_lines = PPTEngine._word_wrap(draw, line_text.strip(), pil_font, max(8, w - text_margin_x * 2))
                    rendered_lines.extend((wl, pil_font, font_color, para.alignment) for wl in wrapped_lines)

                line_height = max(1, max((font.getbbox("Ag")[3] - font.getbbox("Ag")[1] for _, font, _, _ in rendered_lines), default=font_size_px) + 4)
                if getattr(tf, "vertical_anchor", None) == MSO_ANCHOR.MIDDLE:
                    text_y = y + max(text_margin_y, (h - line_height * len(rendered_lines)) // 2)
                elif getattr(tf, "vertical_anchor", None) == MSO_ANCHOR.BOTTOM:
                    text_y = y + max(text_margin_y, h - line_height * len(rendered_lines) - text_margin_y)

                for wl, pil_font, font_color, alignment in rendered_lines:
                    # Determine x based on paragraph alignment
                    text_x = x + text_margin_x
                    try:
                        bbox = draw.textbbox((0, 0), wl, font=pil_font)
                        text_w = bbox[2] - bbox[0]
                        if alignment == PP_ALIGN.CENTER:
                            text_x = x + (w - text_w) // 2
                        elif alignment == PP_ALIGN.RIGHT:
                            text_x = x + w - text_w - 4
                    except Exception:
                        pass

                    if text_y < slide_h_px:
                        draw.text((text_x, text_y), wl, fill=font_color, font=pil_font)
                    text_y += line_height

            # Convert to base64 PNG
            buf = io.BytesIO()
            img.save(buf, format='PNG', quality=92)
            buf.seek(0)
            b64 = base64.b64encode(buf.getvalue()).decode("ascii")
            slide_images.append(b64)

        return slide_images

    @staticmethod
    def _word_wrap(draw, text: str, font, max_width: int) -> list:
        """Word-wrap text to fit within max_width pixels."""
        words = text.split(' ')
        lines = []
        current_line = ""
        for word in words:
            test_line = f"{current_line} {word}".strip()
            try:
                bbox = draw.textbbox((0, 0), test_line, font=font)
                line_w = bbox[2] - bbox[0]
            except Exception:
                line_w = len(test_line) * 8
            if line_w <= max_width:
                current_line = test_line
            else:
                if current_line:
                    lines.append(current_line)
                current_line = word
        if current_line:
            lines.append(current_line)
        return lines if lines else [text]

    @staticmethod
    def _compact_text(value: str) -> str:
        return re.sub(r"[^a-z0-9]+", "", str(value or "").lower())

    @staticmethod
    def _clean_generated_text(value: str) -> str:
        value = re.sub(r"[`*_]", "", str(value or ""))
        value = re.sub(r"\s+", " ", value).strip(" :-")
        return value

    @staticmethod
    def _shape_font_size(shape) -> float:
        sizes = []
        try:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.size:
                        sizes.append(run.font.size.pt)
        except Exception:
            pass
        return max(sizes, default=12.0)

    @staticmethod
    def _is_placeholder_text(text: str) -> bool:
        compact = PPTEngine._compact_text(text)
        return any(token in compact for token in ("loremipsum", "placeholder", "replacewith", "sampletext"))

    @staticmethod
    def _extract_pitch_outline(pitch_sections: Dict[str, Any]) -> list:
        """Extract the optional slide-by-slide part of the streamed pitch."""
        raw = str((pitch_sections or {}).get("full_raw") or "")
        if not raw:
            return []
        outline_match = re.search(r"##\s*Pitch Outline(.*?)(?=##\s*Final Pitch Showcase|\Z)", raw, re.I | re.S)
        outline_text = outline_match.group(1) if outline_match else raw
        records = []
        current = None
        slide_re = re.compile(
            r"^\s*(?:#{1,6}\s*)?(?:\*\*)?slide\s*(\d+)\s*[:.)\-–—]\s*(.+?)(?:\*\*)?\s*$",
            re.I,
        )
        for raw_line in outline_text.splitlines():
            line = raw_line.strip()
            if not line:
                continue
            match = slide_re.match(line)
            if match:
                if current:
                    records.append(current)
                current = {
                    "number": int(match.group(1)),
                    "title": PPTEngine._clean_generated_text(match.group(2)),
                    "bullets": [],
                    "paragraphs": [],
                }
                continue
            if current is None:
                continue
            is_bullet = line.startswith(("- ", "* ", "• "))
            cleaned = PPTEngine._clean_generated_text(line[2:] if is_bullet else line)
            if not cleaned:
                continue
            (current["bullets"] if is_bullet else current["paragraphs"]).append(cleaned)
        if current:
            records.append(current)
        for record in records:
            if not record["bullets"] and record["paragraphs"]:
                record["bullets"] = record["paragraphs"][:4]
            record["bullets"] = [item[:180] for item in record["bullets"][:5]]
        return records

    @staticmethod
    def _slide_role(text: str, index: int) -> str:
        compact = PPTEngine._compact_text(text)
        role_keywords = [
            ("conclusion", ("thankyou", "conclusion", "nextsteps")),
            ("problem", ("problem", "challenge", "painpoint")),
            ("solution", ("solution", "productoverview", "serviceoverview", "uniquevalueproposition")),
            ("market", ("market", "targetaudience", "statistics", "financial", "traction")),
            ("roadmap", ("roadmap", "timeline", "objectives", "milestone")),
            ("team", ("ourteam", "theteam", "superteam", "team", "manager", "staff")),
            ("architecture", ("architecture", "technology", "techstack", "businessmodel")),
            ("intro", ("introduction", "aboutourcompany", "projectoverview", "pitchdeck", "startuppitch")),
        ]
        for role, keywords in role_keywords:
            if any(keyword in compact for keyword in keywords):
                return role
        return "intro" if index == 0 else "general"

    @staticmethod
    def fill_presentation(
        template_source: str,
        session_name: str,
        problem_statement: str,
        user_idea: str,
        pitch_sections: Dict[str, str],
        milestones: List[Dict[str, Any]] = None,
        tasks: List[Dict[str, Any]] = None,
        team_data: Dict[str, Any] = None,
        custom_pptx_bytes: bytes = None
    ) -> bytes:
        """Populate a selected deck while preserving its authored layout and artwork.

        A template is not a fixed ten-slide wireframe: the bundled decks have
        different slide counts and different content roles.  This method therefore
        detects each slide's role from its existing title, replaces only title and
        placeholder copy, and leaves all graphics, pictures, card geometry, and
        footer branding in place.
        """
    @staticmethod
    def _create_default_base_presentation() -> Presentation:
        """Creates a styled, functional 8-slide deck if template files are missing or Git LFS pointers."""
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        bullet_slide_layout = prs.slide_layouts[1]

        roles_config = [
            ("intro", "Project Pitch", "Lorem ipsum dolor sit amet, presentedby : Innovator"),
            ("problem", "The Problem", "Lorem ipsum placeholder problem text.\n• Replacewith first point\n• Sampletext second point"),
            ("solution", "The Solution", "Lorem ipsum placeholder solution text.\n• Replacewith solution point"),
            ("architecture", "How It Works", "Lorem ipsum placeholder architecture text.\n• Replacewith tech stack"),
            ("roadmap", "Roadmap & Execution", "Lorem ipsum placeholder roadmap text.\n• Replacewith milestone"),
            ("market", "Audience & Impact", "Lorem ipsum placeholder market text.\n• Replacewith audience point"),
            ("team", "The Team", "Lorem ipsum placeholder team manager.\n• Replacewith team role"),
            ("conclusion", "Next Steps", "Lorem ipsum placeholder conclusion text.\n• Replacewith next steps")
        ]

        for idx, (role, default_title, default_body) in enumerate(roles_config):
            layout = title_slide_layout if idx == 0 else bullet_slide_layout
            slide = prs.slides.add_slide(layout)
            if slide.shapes.title:
                slide.shapes.title.text = default_title
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = default_body

        return prs

    @staticmethod
    def _load_presentation(t_path: Optional[str] = None, custom_pptx_bytes: Optional[bytes] = None) -> Presentation:
        if custom_pptx_bytes:
            try:
                return Presentation(io.BytesIO(custom_pptx_bytes))
            except Exception as e:
                logger.warning(f"Could not load custom pptx bytes: {e}")

        if t_path and os.path.exists(t_path):
            try:
                if os.path.getsize(t_path) > 5000:
                    return Presentation(t_path)
                else:
                    logger.warning(f"Template at {t_path} is an LFS pointer or invalid (size={os.path.getsize(t_path)}B). Attempting repair.")
            except Exception as e:
                logger.warning(f"Could not open presentation at {t_path}: {e}")

        # Fallback search across TEMPLATES_DIR for any valid .pptx (>5KB)
        if os.path.isdir(TEMPLATES_DIR):
            for fname in os.listdir(TEMPLATES_DIR):
                if fname.endswith(".pptx"):
                    candidate = os.path.join(TEMPLATES_DIR, fname)
                    try:
                        if os.path.getsize(candidate) > 5000:
                            return Presentation(candidate)
                    except Exception:
                        pass

        # Generate default base presentation if no valid PPTX file exists
        default_prs = PPTEngine._create_default_base_presentation()
        if t_path:
            try:
                os.makedirs(os.path.dirname(t_path), exist_ok=True)
                default_prs.save(t_path)
                logger.info(f"Generated and saved base presentation template at {t_path}")
            except Exception as e:
                logger.warning(f"Could not save base presentation template: {e}")
        return default_prs

    @staticmethod
    def fill_presentation(
        template_source: str,
        session_name: str,
        problem_statement: str,
        user_idea: str,
        pitch_sections: Dict[str, str],
        milestones: List[Dict[str, Any]] = None,
        tasks: List[Dict[str, Any]] = None,
        team_data: Dict[str, Any] = None,
        custom_pptx_bytes: bytes = None
    ) -> bytes:
        """Populate a selected deck while preserving its authored layout and artwork.

        A template is not a fixed ten-slide wireframe: the bundled decks have
        different slide counts and different content roles.  This method therefore
        detects each slide's role from its existing title, replaces only title and
        placeholder copy, and leaves all graphics, pictures, card geometry, and
        footer branding in place.
        """
        t_path = PPTEngine.get_template_path(template_source) if not custom_pptx_bytes else None
        prs = PPTEngine._load_presentation(t_path=t_path, custom_pptx_bytes=custom_pptx_bytes)

        team_info = team_data or {}
        members = team_info.get("members") if isinstance(team_info.get("members"), list) else []
        first_member = members[0] if members else team_info
        user_name = first_member.get("name") or first_member.get("full_name") or "Innovator"
        user_role = first_member.get("role") or first_member.get("primary_role") or "Fullstack Engineer"
        skills = first_member.get("skills") or first_member.get("tech_stack") or ["Python", "React"]
        skills_str = ", ".join(map(str, skills)) if isinstance(skills, list) else str(skills)

        milestone_bullets = [
            f"{m.get('phase', 'Milestone')}: {m.get('title') or m.get('deliverable') or 'In progress'}"
            for m in (milestones or [])
        ]
        task_bullets = [f"{t.get('name', 'Task')} · {(t.get('status') or 'pending').replace('_', ' ')}" for t in (tasks or [])[:5]]
        outline = PPTEngine._extract_pitch_outline(pitch_sections or {})
        demo = pitch_sections.get("demo") or "A focused, end-to-end product journey from the user's problem to measurable value."
        showcase = pitch_sections.get("showcase") or "A practical product that turns a real execution problem into a clear, demonstrable outcome."

        content_by_role = {
            "intro": {"title": session_name or "Project Pitch", "summary": user_idea or "A focused product built to solve a real user problem.", "bullets": ["Clear user problem", "Focused product experience", "Measurable execution"]},
            "problem": {"title": "The Problem", "summary": problem_statement or "Teams lose momentum when planning, execution, and communication are disconnected.", "bullets": ["Fragmented planning and execution", "Slow feedback loops", "Too much manual coordination"]},
            "solution": {"title": "The Solution", "summary": demo, "bullets": ["Guided product workflow", "Real-time progress visibility", "Actionable AI assistance"]},
            "architecture": {"title": "How It Works", "summary": pitch_sections.get("architecture") or "A lightweight web application connects the product experience, data layer, and AI workflow.", "bullets": ["FastAPI service layer", "React product experience", "Persistent project data and AI orchestration"]},
            "roadmap": {"title": "Roadmap & Execution", "summary": "The build is organized into small, demonstrable milestones.", "bullets": milestone_bullets or ["Validate the core workflow", "Ship the product experience", "Measure and iterate"]},
            "market": {"title": "Audience & Impact", "summary": "The product is designed around the people who need a faster, clearer path from idea to outcome.", "bullets": ["Primary users with an urgent workflow", "A repeatable product moment", "Clear path to measurable impact"]},
            "team": {"title": "The Team", "summary": f"{user_name} leads the build as {user_role}.", "bullets": [f"Lead: {user_name}", f"Role: {user_role}", f"Skills: {skills_str}"]},
            "conclusion": {"title": "Next Steps", "summary": showcase, "bullets": ["Demo the core workflow", "Show the measured outcome", "Invite the next step"]},
            "general": {"title": "Project Overview", "summary": showcase, "bullets": task_bullets or ["Product capability", "Technical foundation", "Execution outcome"]},
        }

        def choose_content(role, slide_index):
            content = dict(content_by_role.get(role, content_by_role["general"]))
            matching = None
            role_tokens = {
                "problem": ("problem",),
                "solution": ("solution",),
                "market": ("market", "traction", "statistics", "financial"),
                "roadmap": ("roadmap", "timeline", "objective"),
                "team": ("team",),
                "architecture": ("architecture", "technology", "businessmodel"),
                "conclusion": ("conclusion", "thankyou"),
                "intro": ("intro", "overview", "pitch"),
            }
            for item in outline:
                item_text = PPTEngine._compact_text(item.get("title"))
                if role != "general" and any(token in item_text for token in role_tokens.get(role, (role,))):
                    matching = item
                    break
            if matching is None and slide_index < len(outline):
                matching = outline[slide_index]
            if matching:
                content["title"] = matching.get("title") or content["title"]
                content["bullets"] = matching.get("bullets") or content["bullets"]
                if matching.get("paragraphs"):
                    content["summary"] = matching["paragraphs"][0]
            return content

        footer_tokens = ("ingoudecompany", "fradelandspies", "thynkunlimited", "hello", "www", "123anywhere", "august", "december", "website")
        title_tokens = ("pitchdeck", "problem", "solution", "roadmap", "introduction", "overview", "market", "team", "thankyou", "conclusion", "objectives", "timeline", "traction", "businessmodel", "architecture", "statistics", "product")

        for slide_index, slide in enumerate(prs.slides):
            text_shapes = [shape for shape in slide.shapes if shape.has_text_frame and shape.text_frame.text.strip()]
            if not text_shapes:
                continue

            slide_text = " ".join(shape.text_frame.text for shape in text_shapes)
            role = PPTEngine._slide_role(slide_text, slide_index)
            content = choose_content(role, slide_index)

            # Determine the authored title from the template, not from a fixed
            # slide number.  This works across all five bundled decks and most
            # user-uploaded decks that use ordinary title text boxes.
            title_candidates = [
                shape for shape in text_shapes
                if any(token in PPTEngine._compact_text(shape.text_frame.text) for token in title_tokens)
                and not any(token in PPTEngine._compact_text(shape.text_frame.text) for token in footer_tokens)
            ]
            title_shape = max(title_candidates, key=PPTEngine._shape_font_size) if title_candidates else None
            if title_shape is None and slide_index == 0:
                title_shape = max(text_shapes, key=PPTEngine._shape_font_size)
            placeholder_shapes = [candidate for candidate in text_shapes if PPTEngine._is_placeholder_text(candidate.text_frame.text)]

            for shape in text_shapes:
                original = shape.text_frame.text.strip()
                compact = PPTEngine._compact_text(original)
                if "presentedby" in compact:
                    PPTEngine.fit_text_to_frame(shape.text_frame, f"Presented By : {user_name}", max_font_size=12, min_font_size=8)
                    continue
                if shape is title_shape:
                    PPTEngine.fit_text_to_frame(shape.text_frame, content["title"], max_font_size=32, min_font_size=14, is_title=True)
                    continue
                if any(token in compact for token in footer_tokens) or compact in {"001", "002", "003", "1", "2", "3"}:
                    continue
                if any(token in compact for token in ("manager", "staf", "projectmanager")) and role == "team":
                    PPTEngine.fit_text_to_frame(shape.text_frame, user_role, max_font_size=12, min_font_size=8)
                    continue
                if not PPTEngine._is_placeholder_text(original):
                    continue

                # Keep every authored card in its original location.  Feed the
                # summary to the first large placeholder and one short bullet to
                # each remaining card, which avoids overflowing a card with a
                # whole paragraph of LLM output.
                placeholder_index = placeholder_shapes.index(shape)
                if len(placeholder_shapes) == 1:
                    replacement = "\n".join([content["summary"]] + [f"• {bullet}" for bullet in content["bullets"][:3]])
                elif placeholder_index == 0:
                    replacement = content["summary"]
                else:
                    bullets = content["bullets"] or [content["summary"]]
                    bullet_index = placeholder_index - 1
                    replacement = bullets[bullet_index] if bullet_index < len(bullets) else ""
                PPTEngine.fit_text_to_frame(shape.text_frame, replacement, max_font_size=18, min_font_size=8)

        output_stream = io.BytesIO()
        prs.save(output_stream)
        output_stream.seek(0)
        return output_stream.getvalue()

    @staticmethod
    def generate_project_pdf(
        session_name: str,
        problem_statement: str = "",
        user_idea: str = "",
        milestones: List[Dict[str, Any]] = None,
        tasks: List[Dict[str, Any]] = None,
        blockers: List[Dict[str, Any]] = None,
        team_data: Dict[str, Any] = None,
        pitch_sections: Dict[str, str] = None,
        created_at: str = None
    ) -> bytes:
        milestones = milestones or []
        tasks = tasks or []
        blockers = blockers or []
        team_data = team_data or {}
        pitch_sections = pitch_sections or {}

        buffer = io.BytesIO()
        doc = SimpleDocTemplate(
            buffer,
            pagesize=letter,
            rightMargin=36,
            leftMargin=36,
            topMargin=54,
            bottomMargin=54
        )

        styles = getSampleStyleSheet()
        
        # Clean Hackathon Theme Typography
        title_style = ParagraphStyle(
            'HackTitle',
            parent=styles['Heading1'],
            fontName='Helvetica-Bold',
            fontSize=18,
            leading=22,
            textColor=colors.HexColor('#ffffff')
        )

        subtitle_style = ParagraphStyle(
            'HackSubTitle',
            parent=styles['Normal'],
            fontName='Helvetica-Bold',
            fontSize=9,
            leading=12,
            textColor=colors.HexColor('#c084fc'),
            alignment=2
        )

        section_title_style = ParagraphStyle(
            'SectionTitle',
            parent=styles['Heading2'],
            fontName='Helvetica-Bold',
            fontSize=13,
            leading=17,
            textColor=colors.HexColor('#6d28d9'),
            spaceBefore=14,
            spaceAfter=6,
            keepWithNext=True
        )

        body_style = ParagraphStyle(
            'ReportBody',
            parent=styles['BodyText'],
            fontName='Helvetica',
            fontSize=9.5,
            leading=14,
            textColor=colors.HexColor('#334155'),
            spaceAfter=6
        )

        callout_style = ParagraphStyle(
            'ReportCallout',
            parent=body_style,
            fontName='Helvetica',
            fontSize=9.5,
            leading=14,
            textColor=colors.HexColor('#1e293b'),
            backColor=colors.HexColor('#f8fafc'),
            borderColor=colors.HexColor('#7c3aed'),
            borderWidth=1,
            borderPadding=10,
            spaceAfter=10
        )

        table_cell_style = ParagraphStyle(
            'TableCell',
            parent=styles['Normal'],
            fontName='Helvetica',
            fontSize=8.5,
            leading=11.5,
            textColor=colors.HexColor('#1e293b')
        )

        table_header_style = ParagraphStyle(
            'TableHeader',
            parent=styles['Normal'],
            fontName='Helvetica-Bold',
            fontSize=8.5,
            leading=11.5,
            textColor=colors.white
        )

        story = []

        # Document Header Banner - Modern Purple Hackathon Header
        header_table = Table(
            [[Paragraph(f"<b>{session_name}</b>", title_style),
              Paragraph(f"<b>KAIROS AI EXECUTION DOSSIER</b><br/><font color='#e9d5ff'>Hackathon Blueprint Report</font>", subtitle_style)]],
            colWidths=[330, 210]
        )
        header_table.setStyle(TableStyle([
            ('BACKGROUND', (0,0), (-1,-1), colors.HexColor('#6d28d9')),
            ('PADDING', (0,0), (-1,-1), 12),
            ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
        ]))
        story.append(header_table)
        story.append(Spacer(1, 14))

        # KPI Metrics Cards Table
        total_tasks = len(tasks)
        completed_tasks = sum(1 for t in tasks if t.get('status') == 'completed')
        pct = int((completed_tasks / total_tasks * 100)) if total_tasks > 0 else 0
        open_blockers = sum(1 for b in blockers if b.get('status') == 'open')

        kpi_data = [
            [
                Paragraph(f"<font size=7 color='#64748b'><b>MILESTONES</b></font><br/><font size=15 color='#6d28d9'><b>{len(milestones)}</b></font>", table_cell_style),
                Paragraph(f"<font size=7 color='#64748b'><b>COMPLETED TASKS</b></font><br/><font size=15 color='#059669'><b>{completed_tasks}/{total_tasks}</b></font>", table_cell_style),
                Paragraph(f"<font size=7 color='#64748b'><b>PROGRESS</b></font><br/><font size=15 color='#2563eb'><b>{pct}%</b></font>", table_cell_style),
                Paragraph(f"<font size=7 color='#64748b'><b>OPEN BLOCKERS</b></font><br/><font size=15 color='#dc2626'><b>{open_blockers}</b></font>", table_cell_style)
            ]
        ]
        kpi_table = Table(kpi_data, colWidths=[135, 135, 135, 135])
        kpi_table.setStyle(TableStyle([
            ('BACKGROUND', (0,0), (-1,-1), colors.HexColor('#f8fafc')),
            ('BOX', (0,0), (-1,-1), 1, colors.HexColor('#e2e8f0')),
            ('INNERGRID', (0,0), (-1,-1), 0.5, colors.HexColor('#cbd5e1')),
            ('PADDING', (0,0), (-1,-1), 6),
            ('ALIGN', (0,0), (-1,-1), 'CENTER'),
        ]))
        story.append(kpi_table)
        story.append(Spacer(1, 14))

        # 1. Project Overview & Problem Statement
        story.append(Paragraph("1. Project Overview & Problem Statement", section_title_style))
        story.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor("#7c3aed"), spaceAfter=8))

        if problem_statement:
            story.append(Paragraph(f"<b>Problem Statement:</b><br/>{problem_statement}", callout_style))
        
        if user_idea:
            story.append(Paragraph(f"<b>Proposed Solution & Architecture:</b><br/>{user_idea}", body_style))

        # 2. Team Roster & Skill Matrix
        if team_data:
            story.append(Spacer(1, 8))
            story.append(Paragraph("2. Team Structure & Skill Alignment", section_title_style))
            story.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor("#7c3aed"), spaceAfter=6))
            
            members = team_data.get("members", [])
            if not members and "name" in team_data:
                members = [team_data]
                
            if members:
                t_rows = [[
                    Paragraph("<b>Member Name</b>", table_header_style),
                    Paragraph("<b>Role</b>", table_header_style),
                    Paragraph("<b>Level</b>", table_header_style),
                    Paragraph("<b>Tech Stack Skills</b>", table_header_style)
                ]]
                for m in members:
                    skills_str = ", ".join(m.get("skills", [])) if isinstance(m.get("skills"), list) else str(m.get("skills", "N/A"))
                    t_rows.append([
                        Paragraph(f"<b>{m.get('full_name') or m.get('name') or 'Developer'}</b>", table_cell_style),
                        Paragraph(m.get("role") or m.get("primary_role") or "Developer", table_cell_style),
                        Paragraph(m.get("level") or m.get("experience_level") or "Mid", table_cell_style),
                        Paragraph(skills_str or "Fullstack", table_cell_style)
                    ])
                t_team = Table(t_rows, colWidths=[130, 110, 80, 220])
                t_team.setStyle(TableStyle([
                    ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#1e1b4b')),
                    ('PADDING', (0,0), (-1,-1), 5),
                    ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#e2e8f0')),
                    ('ROWBACKGROUNDS', (0,1), (-1,-1), [colors.HexColor('#ffffff'), colors.HexColor('#f8fafc')])
                ]))
                story.append(t_team)

        # 3. Strategic Roadmap & Execution Milestones
        story.append(Spacer(1, 10))
        story.append(Paragraph("3. Strategic Roadmap & Execution Milestones", section_title_style))
        story.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor("#7c3aed"), spaceAfter=6))

        if milestones:
            m_rows = [[
                Paragraph("<b>Phase</b>", table_header_style),
                Paragraph("<b>Milestone Title</b>", table_header_style),
                Paragraph("<b>Key Deliverable</b>", table_header_style),
                Paragraph("<b>Est. Duration</b>", table_header_style)
            ]]
            for m in milestones:
                m_rows.append([
                    Paragraph(f"<b>{m.get('phase', 'Phase')}</b>", table_cell_style),
                    Paragraph(f"<b>{m.get('title', '')}</b>", table_cell_style),
                    Paragraph(m.get('deliverable', 'Code & Architecture'), table_cell_style),
                    Paragraph(m.get('duration_estimate', '2-4 hrs'), table_cell_style)
                ])
            m_table = Table(m_rows, colWidths=[80, 180, 200, 80])
            m_table.setStyle(TableStyle([
                ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#6d28d9')),
                ('PADDING', (0,0), (-1,-1), 5),
                ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#e2e8f0')),
                ('ROWBACKGROUNDS', (0,1), (-1,-1), [colors.HexColor('#ffffff'), colors.HexColor('#f8fafc')])
            ]))
            story.append(m_table)

        # 4. Comprehensive Task Execution Matrix
        story.append(Spacer(1, 10))
        story.append(Paragraph("4. Task Execution Matrix", section_title_style))
        story.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor("#7c3aed"), spaceAfter=6))

        if tasks:
            task_rows = [[
                Paragraph("<b>Task Name</b>", table_header_style),
                Paragraph("<b>Priority</b>", table_header_style),
                Paragraph("<b>Status</b>", table_header_style)
            ]]
            for idx, t in enumerate(tasks):
                status_val = (t.get('status') or 'pending').lower()
                status_color = '#059669' if status_val == 'completed' else '#dc2626' if status_val == 'blocked' else '#d97706' if status_val == 'in_progress' else '#0284c7'
                status_text = f"<font color='{status_color}'><b>{status_val.capitalize()}</b></font>"

                prio_val = (t.get('priority') or 'medium').upper()
                prio_color = '#dc2626' if prio_val == 'HIGH' else '#7c3aed' if prio_val == 'MEDIUM' else '#0284c7'
                prio_text = f"<font color='{prio_color}'><b>{prio_val}</b></font>"

                task_rows.append([
                    Paragraph(t.get('name', 'Task'), table_cell_style),
                    Paragraph(prio_text, table_cell_style),
                    Paragraph(status_text, table_cell_style)
                ])
            task_table = Table(task_rows, colWidths=[340, 100, 100])
            task_table.setStyle(TableStyle([
                ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#1e1b4b')),
                ('PADDING', (0,0), (-1,-1), 4.5),
                ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#e2e8f0')),
                ('ROWBACKGROUNDS', (0,1), (-1,-1), [colors.HexColor('#ffffff'), colors.HexColor('#f8fafc')])
            ]))
            story.append(task_table)

        # 5. Technical Blockers & Challenges
        if blockers:
            story.append(Spacer(1, 10))
            story.append(Paragraph("5. Technical Challenges & Blockers", section_title_style))
            story.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor("#7c3aed"), spaceAfter=6))

            b_rows = [[
                Paragraph("<b>Blocker Description</b>", table_header_style),
                Paragraph("<b>Severity</b>", table_header_style),
                Paragraph("<b>Status</b>", table_header_style)
            ]]
            for b in blockers:
                b_rows.append([
                    Paragraph(b.get('description', ''), table_cell_style),
                    Paragraph(f"<b>{b.get('severity', 'medium').upper()}</b>", table_cell_style),
                    Paragraph(f"<b>{b.get('status', 'open').capitalize()}</b>", table_cell_style)
                ])
            b_table = Table(b_rows, colWidths=[360, 90, 90])
            b_table.setStyle(TableStyle([
                ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#991b1b')),
                ('PADDING', (0,0), (-1,-1), 5),
                ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#fee2e2')),
            ]))
            story.append(b_table)

        # 6. Presentation Suite & Pitch Script
        raw_pitch = pitch_sections.get("full_raw") or pitch_sections.get("showcase") or ""
        if raw_pitch:
            story.append(PageBreak())
            story.append(Paragraph("6. Presentation Suite & Stage Pitch Script", section_title_style))
            story.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor("#7c3aed"), spaceAfter=8))

            lines = raw_pitch.split("\n")
            for line in lines:
                l_str = line.strip()
                if not l_str:
                    story.append(Spacer(1, 3))
                elif l_str.startswith("#"):
                    clean_h = re.sub(r'#+\s*', '', l_str)
                    story.append(Paragraph(f"<b>{clean_h}</b>", ParagraphStyle('PitchH', parent=section_title_style, fontSize=11, leading=15, spaceBefore=6, spaceAfter=3)))
                elif l_str.startswith("-") or l_str.startswith("*"):
                    clean_bullet = l_str[1:].strip()
                    story.append(Paragraph(f"• {clean_bullet}", ParagraphStyle('PitchBullet', parent=body_style, leftIndent=12, firstLineIndent=-8)))
                else:
                    story.append(Paragraph(l_str, body_style))

        doc.build(story, canvasmaker=NumberedCanvas)
        buffer.seek(0)
        return buffer.getvalue()

    @staticmethod
    def generate_pdf(
        session_name: str,
        problem_statement: str = "",
        user_idea: str = "",
        pitch_sections: Dict[str, str] = None,
        milestones: List[Dict[str, Any]] = None,
        tasks: List[Dict[str, Any]] = None
    ) -> bytes:
        return PPTEngine.generate_project_pdf(
            session_name=session_name,
            problem_statement=problem_statement,
            user_idea=user_idea,
            milestones=milestones,
            tasks=tasks,
            pitch_sections=pitch_sections
        )
